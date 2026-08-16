from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE


OUT = "Capabl_Job_Search_AI_Agent_Final.pptx"

W, H = 13.333, 7.5

NAVY = "0B1020"
NAVY_2 = "11182B"
INK = "111827"
SLATE = "475569"
MUTED = "64748B"
LINE = "DCE3EE"
PALE = "F5F7FB"
WHITE = "FFFFFF"
INDIGO = "4F46E5"
BLUE = "2563EB"
CYAN = "06B6D4"
GREEN = "10B981"
AMBER = "F59E0B"
RED = "EF4444"
VIOLET = "7C3AED"
YELLOW = "FFD400"


def rgb(hex_value):
    return RGBColor.from_string(hex_value)


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False,
             font="Aptos", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP,
             margin=0, italic=False, linespacing=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = linespacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = rgb(color)
    return box


def add_rich_text(slide, runs, x, y, w, h, size=18, color=INK,
                  valign=MSO_ANCHOR.TOP, margin=0, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear(); tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.alignment = align
    for item in runs:
        r = p.add_run(); r.text = item[0]
        r.font.name = "Aptos"; r.font.size = Pt(item[1] if len(item) > 1 else size)
        r.font.bold = item[2] if len(item) > 2 else False
        r.font.color.rgb = rgb(item[3] if len(item) > 3 else color)
    return box


def shape(slide, kind, x, y, w, h, fill=WHITE, line=LINE, radius=True, lw=1):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else kind,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
    s.line.color.rgb = rgb(line); s.line.width = Pt(lw)
    return s


def rect(slide, x, y, w, h, fill, line=None, radius=False, lw=0):
    return shape(slide, MSO_SHAPE.RECTANGLE, x, y, w, h, fill,
                 fill if line is None else line, radius, lw)


def circle(slide, x, y, d, fill, line=None, lw=0):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = rgb(fill)
    s.line.color.rgb = rgb(fill if line is None else line); s.line.width = Pt(lw)
    return s


def line(slide, x1, y1, x2, y2, color=LINE, width=1.5, dash=None):
    s = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    s.line.color.rgb = rgb(color); s.line.width = Pt(width)
    if dash: s.line.dash_style = dash
    return s


def pill(slide, text, x, y, w, fill="EEF2FF", color=INDIGO, border=None, size=10):
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, .31, fill, border or fill, True, .7)
    add_text(slide, text, x, y+.005, w, .29, size, color, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def icon_badge(slide, label, x, y, fill=INDIGO, d=.42, size=13):
    circle(slide, x, y, d, fill)
    add_text(slide, label, x, y+.005, d, d-.01, size, WHITE, True,
             align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def capabl_logo(slide, x=.55, y=.34, dark=False):
    c = WHITE if dark else INK
    add_text(slide, "Capabl", x, y, .72, .25, 12, c, True)
    rect(slide, x+.67, y+.055, .11, .11, YELLOW, YELLOW)


def footer(slide, n, dark=False):
    color = "98A2B3" if dark else MUTED
    add_text(slide, "TEAM QUADCORE  •  JOB SEARCH AI AGENT", .55, 7.15, 4.2, .18, 8, color, True)
    add_text(slide, f"{n:02d}", 12.35, 7.12, .42, .2, 9, color, True, align=PP_ALIGN.RIGHT)


def title_block(slide, eyebrow, title, subtitle=None, dark=False):
    c = WHITE if dark else INK
    pill(slide, eyebrow.upper(), .62, .48, 1.55, fill=("222A42" if dark else "EEF2FF"), color=("A5B4FC" if dark else INDIGO))
    add_text(slide, title, .62, .95, 12.0, .58, 27, c, True)
    if subtitle:
        add_text(slide, subtitle, .64, 1.53, 11.7, .40, 12.2, ("A7B0C4" if dark else MUTED))


def add_bullet(slide, text, x, y, w, color=SLATE, dot=INDIGO, size=12, bold_lead=None):
    circle(slide, x, y+.12, .08, dot)
    if bold_lead and text.startswith(bold_lead):
        add_rich_text(slide, [(bold_lead, size, True, INK), (text[len(bold_lead):], size, False, color)], x+.18, y, w-.18, .45)
    else:
        add_text(slide, text, x+.18, y, w-.18, .45, size, color)


prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]
prs.core_properties.title = "Job Search AI Agent — Team QuadCore"
prs.core_properties.subject = "Capabl project presentation updated from repository implementation"
prs.core_properties.author = "Team QuadCore"
prs.core_properties.comments = "Generated as an editable 16:9 deck from the current repository state."


# 1 — Cover
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, NAVY, NAVY)
circle(slide, 9.55, -.8, 4.4, "1E2A55")
circle(slide, 10.55, .15, 2.5, INDIGO)
circle(slide, 11.28, .88, 1.05, CYAN)
capabl_logo(slide, .62, .44, dark=True)
pill(slide, "TEAM QUADCORE", .62, 1.35, 1.65, fill="222A42", color="A5B4FC")
add_text(slide, "Job Search\nAI Agent", .62, 1.88, 7.4, 1.8, 39, WHITE, True)
add_text(slide, "A resume-aware career workspace for discovery, decisions and interview readiness.", .66, 3.86, 6.7, .68, 17, "B7C0D8")

# Decorative agent path
line(slide, 8.4, 5.0, 11.9, 5.0, "4B5C89", 2)
for i, (x, label, col) in enumerate([(8.4,"CV",BLUE),(9.55,"MATCH",VIOLET),(10.85,"ACT",CYAN),(11.9,"READY",GREEN)]):
    circle(slide, x-.19, 4.81, .38, col)
    add_text(slide, str(i+1), x-.19, 4.815, .38, .36, 10, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, label, x-.45, 5.30, .9, .2, 8.5, "A7B0C4", True, align=PP_ALIGN.CENTER)

shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, .62, 5.45, 6.3, 1.17, "121B31", "2A3553", True, 1)
add_text(slide, "Gangadhar R  •  Gururaj S H  •  H M Akshay  •  Mallikarjunayya S", .92, 5.78, 5.72, .28, 13, WHITE, True, align=PP_ALIGN.CENTER)
add_text(slide, "Capabl project presentation  •  2-month build", .92, 6.12, 5.72, .2, 10.5, "8792AD", align=PP_ALIGN.CENTER)
footer(slide, 1, dark=True)


# 2 — Challenge
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, WHITE, WHITE)
capabl_logo(slide)
title_block(slide, "The challenge", "The job search is a broken chain of context",
            "Candidates repeatedly translate the same experience across disconnected tools.")

labels = [
    ("01", "Resume", "Static PDF\nLow feedback", BLUE),
    ("02", "Discovery", "Generic lists\nManual filtering", VIOLET),
    ("03", "Application", "Repeated tailoring\nScattered tracking", AMBER),
    ("04", "Interview", "Generic practice\nLate feedback", GREEN),
]
for i, (num, head, body, col) in enumerate(labels):
    x = .68 + i*3.12
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.25, 2.76, 2.03, PALE, LINE, True, 1)
    icon_badge(slide, num, x+.22, 2.48, col, .46, 10)
    add_text(slide, head, x+.22, 3.06, 2.25, .30, 16, INK, True)
    add_text(slide, body, x+.22, 3.48, 2.22, .52, 11.5, MUTED)
    if i < 3:
        add_text(slide, "→", x+2.82, 3.06, .28, .28, 20, "A8B1C1", True, align=PP_ALIGN.CENTER)

shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, .68, 4.72, 12.0, 1.54, "EEF2FF", "C7D2FE", True, 1)
add_text(slide, "PRODUCT THESIS", .98, 5.0, 1.55, .2, 9.5, INDIGO, True)
add_text(slide, "Keep the resume as shared context—then use it to rank opportunities, guide applications and personalize preparation.", .98, 5.32, 10.9, .52, 19, INK, True)
footer(slide, 2)


# 3 — Product experience
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, WHITE, WHITE)
capabl_logo(slide)
title_block(slide, "Product overview", "One workspace. Four connected decisions.",
            "The React experience turns a resume into a continuous career workflow.")

steps = [
    ("1", "UNDERSTAND", "Upload a PDF", "Extract skills, experience, education and an ATS health score.", BLUE),
    ("2", "DISCOVER", "Find better-fit roles", "Search live listings and rank results against the selected resume.", VIOLET),
    ("3", "ACT", "Move applications", "Save jobs, add notes and move them through a visual pipeline.", AMBER),
    ("4", "PREPARE", "Practice deliberately", "Generate role-specific questions and receive scored feedback.", GREEN),
]
for i, (n, eyebrow, head, body, col) in enumerate(steps):
    x = .68 + i*3.12
    shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, 2.28, 2.72, 3.62, WHITE, LINE, True, 1)
    rect(slide, x, 2.28, 2.72, .12, col, col, radius=True)
    icon_badge(slide, n, x+.22, 2.67, col, .55, 14)
    add_text(slide, eyebrow, x+.22, 3.48, 2.15, .18, 8.5, col, True)
    add_text(slide, head, x+.22, 3.82, 2.15, .56, 17, INK, True)
    add_text(slide, body, x+.22, 4.55, 2.18, .83, 11.2, MUTED)
    pill(slide, ["PDF → PROFILE", "QUERY → RANK", "SAVE → STATUS", "ANSWER → COACH"][i], x+.22, 5.38, 1.73, fill=PALE, color=SLATE, size=8.2)
footer(slide, 3)


# 4 — Resume intelligence
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, PALE, PALE)
capabl_logo(slide)
title_block(slide, "Resume intelligence", "From PDF to a usable career profile",
            "Deterministic extraction establishes the profile; Gemini fills the narrative gap when needed.")

# UI mockup
shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, .68, 2.18, 6.25, 4.36, NAVY, NAVY, True, 0)
add_text(slide, "Resume Analysis", 1.02, 2.49, 2.5, .28, 14, WHITE, True)
pill(slide, "PDF UPLOADED", 5.1, 2.46, 1.42, fill="1C2947", color="93C5FD", size=8.7)
shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.02, 3.03, 5.56, 1.05, "171F34", "2A3550", True, 1)
icon_badge(slide, "CV", 1.25, 3.31, BLUE, .48, 10)
add_text(slide, "akshay_resume.pdf", 1.91, 3.27, 2.7, .23, 12, WHITE, True)
add_text(slide, "ATS HEALTH", 4.87, 3.23, 1.15, .18, 8, "8490AA", True, align=PP_ALIGN.RIGHT)
add_text(slide, "78", 5.35, 3.49, .68, .32, 21, GREEN, True, align=PP_ALIGN.RIGHT)
add_text(slide, "Extracted skills", 1.05, 4.38, 2.0, .2, 9, "8490AA", True)
for idx, (txt, wid) in enumerate([("Python",.83),("FastAPI",.95),("React",.78),("PostgreSQL",1.22),("Docker",.82)]):
    x = 1.05 + sum([.93,1.05,.88,1.32,.92][:idx])
    pill(slide, txt, x, 4.72, wid, fill="202A45", color="B8C0FF", size=8.8)
add_text(slide, "Profile summary", 1.05, 5.28, 2.0, .2, 9, "8490AA", True)
add_text(slide, "Backend-focused engineer with practical experience building API-led products and data workflows.", 1.05, 5.59, 5.1, .55, 11, "D5DBE9")

# Explanation side
features = [
    ("01", "Text extraction", "pdfplumber reads text from PDF resumes; uploads are capped at 5 MB.", BLUE),
    ("02", "Structured signals", "Rules identify known skills, education, companies and years of experience.", VIOLET),
    ("03", "ATS health", "A transparent heuristic scores skills density and profile completeness.", AMBER),
    ("04", "Stored context", "Parsed data, raw text and a TF‑IDF vector are persisted with the resume.", GREEN),
]
for i, (n, head, body, col) in enumerate(features):
    y = 2.18 + i*1.05
    icon_badge(slide, n, 7.36, y+.03, col, .38, 8.5)
    add_text(slide, head, 7.9, y, 2.1, .24, 13, INK, True)
    add_text(slide, body, 7.9, y+.31, 4.3, .47, 10.5, MUTED)
shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.36, 6.11, 5.26, .45, "FFF7E6", "FDE2A7", True, 1)
add_text(slide, "Note: the ATS score is guidance—not a prediction of any employer’s ATS.", 7.58, 6.235, 4.86, .19, 8.8, "8A5A00", True)
footer(slide, 4)


# 5 — Match and advice
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, WHITE, WHITE)
capabl_logo(slide)
title_block(slide, "Discovery + guidance", "Rank the search. Explain the fit.",
            "Search results become decisions when relevance and next actions are visible together.")

# Left pipeline
shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, .68, 2.17, 4.14, 4.24, NAVY, NAVY, True, 0)
add_text(slide, "MATCH PIPELINE", 1.0, 2.50, 2.0, .2, 9, "8EA0C8", True)
pipe = [
    ("Q", "Search", "Role + location", BLUE),
    ("A", "Retrieve", "Adzuna / mock fallback", CYAN),
    ("↗", "Rank", "TF‑IDF + cosine", VIOLET),
    ("✓", "Explain", "Matched skills + gaps", GREEN),
]
for i, (ico, head, body, col) in enumerate(pipe):
    y=2.98+i*.78
    icon_badge(slide, ico, 1.0, y, col, .40, 10)
    add_text(slide, head, 1.58, y-.01, 1.05, .20, 11.5, WHITE, True)
    add_text(slide, body, 2.66, y-.01, 1.72, .26, 9.5, "A6B0C7")
    if i<3: line(slide, 1.20, y+.41, 1.20, y+.75, "3A4665", 1.5)
pill(slide, "REDIS CACHE  •  TTL 1 HOUR", 1.0, 6.0, 2.58, fill="1B2744", color="A5B4FC", size=8.4)

# Right job/advice mockup
shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 5.17, 2.17, 7.44, 4.24, PALE, LINE, True, 1)
add_text(slide, "Backend Engineer", 5.55, 2.55, 3.4, .30, 18, INK, True)
add_text(slide, "Razorpay  •  Bengaluru", 5.55, 2.94, 3.0, .21, 10.5, MUTED)
circle(slide, 10.95, 2.46, .90, GREEN)
add_text(slide, "82", 10.95, 2.52, .90, .30, 22, WHITE, True, align=PP_ALIGN.CENTER)
add_text(slide, "MATCH", 10.95, 2.83, .90, .17, 7.2, WHITE, True, align=PP_ALIGN.CENTER)
line(slide, 5.55, 3.42, 12.18, 3.42, LINE, 1)
add_text(slide, "MATCHED", 5.55, 3.74, 1.15, .18, 8.4, GREEN, True)
for i, t in enumerate(["Python", "FastAPI", "REST"]): pill(slide, t, 5.55+i*.92, 4.05, .80, fill="E7F8F1", color="087A5B", size=8.5)
add_text(slide, "SKILL GAPS", 8.55, 3.74, 1.15, .18, 8.4, AMBER, True)
for i, t in enumerate(["AWS", "Kubernetes"]): pill(slide, t, 8.55+i*.90, 4.05, .80 if i==0 else 1.18, fill="FFF5DF", color="9A6200", size=8.5)
add_text(slide, "AI NEXT ACTIONS", 5.55, 4.72, 1.75, .18, 8.4, INDIGO, True)
add_bullet(slide, "Tailored improvement suggestions", 5.56, 5.02, 3.1, size=10.5)
add_bullet(slide, "Cover-letter draft for the role", 5.56, 5.42, 3.1, size=10.5)
add_bullet(slide, "Role-specific interview tips", 8.87, 5.02, 2.9, size=10.5)
add_bullet(slide, "Graceful fallback if Gemini fails", 8.87, 5.42, 3.0, size=10.5)
footer(slide, 5)


# 6 — Application cockpit and interview
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, PALE, PALE)
capabl_logo(slide)
title_block(slide, "Act + prepare", "Turn matches into momentum",
            "The workspace continues after discovery: track every application and rehearse for the next conversation.")

# Kanban panel
shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, .68, 2.16, 6.0, 4.35, WHITE, LINE, True, 1)
add_text(slide, "APPLICATION PIPELINE", 1.0, 2.48, 2.4, .2, 9, SLATE, True)
cols=[("Saved",BLUE,2),("Applied",VIOLET,1),("Interviewing",AMBER,1),("Offered",GREEN,1)]
for i,(label,col,count) in enumerate(cols):
    x=.98+i*1.36
    rect(slide,x,2.91,1.17,.08,col,col,radius=True)
    add_text(slide,label,x,3.16,1.17,.38,9.6,INK,True,align=PP_ALIGN.CENTER)
    for j in range(count):
        shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,x,3.72+j*.78,1.17,.61,PALE,LINE,True,1)
        add_text(slide,["Data Analyst","API Engineer","Product Dev","Backend Eng.","SWE I"][i+j],x+.10,3.83+j*.78,.97,.18,8.3,SLATE,True,align=PP_ALIGN.CENTER)
add_text(slide,"Status updates  •  Notes  •  Match score  •  Direct job link",1.0,6.06,5.35,.18,9.2,MUTED,align=PP_ALIGN.CENTER)

# Interview panel
shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 7.0, 2.16, 5.62, 4.35, NAVY, NAVY, True, 0)
add_text(slide, "MOCK INTERVIEW", 7.34, 2.48, 2.2, .2, 9, "8EA0C8", True)
add_text(slide, "Question 3 of 5", 7.34, 2.90, 2.0, .27, 17, WHITE, True)
pill(slide, "TECHNICAL", 10.75, 2.84, 1.28, fill="202B48", color="A5B4FC", size=8.3)
add_text(slide, "How would you design a resilient job-search API with cached results?", 7.34, 3.42, 4.75, .67, 15, WHITE, True)
shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,7.34,4.30,4.92,.78,"151E34","2A3550",True,1)
add_text(slide,"Answer is evaluated against a role-aware framework.",7.58,4.54,4.42,.26,10,"BEC7DA")
add_text(slide,"8.4",7.34,5.46,.82,.37,24,GREEN,True)
add_text(slide,"/ 10",8.13,5.57,.55,.20,10,"8290AA",True)
add_text(slide,"Strengths",9.0,5.37,.85,.18,8.3,GREEN,True)
add_text(slide,"Clear trade-offs",9.0,5.66,1.17,.18,9,"D6DCEA")
add_text(slide,"Improve",10.55,5.37,.75,.18,8.3,AMBER,True)
add_text(slide,"Add metrics",10.55,5.66,1.25,.18,9,"D6DCEA")
add_text(slide,"Technical  •  Behavioral  •  Situational  •  HR",7.34,6.10,4.83,.18,8.7,"8792AD",align=PP_ALIGN.CENTER)
footer(slide, 6)


# 7 — Agent workflow
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, NAVY, NAVY)
capabl_logo(slide, dark=True)
title_block(slide, "Agent workflow", "A stateful four-node career agent", "LangGraph makes each step explicit, testable and observable.", dark=True)

nodes=[
    ("01","PARSE RESUME","raw text\nprofile + ATS\nTF‑IDF vector",BLUE),
    ("02","SEARCH JOBS","Adzuna API\nRedis cache\nmock fallback",CYAN),
    ("03","RANK MATCHES","cosine score\nmatched skills\nskill gaps",VIOLET),
    ("04","GENERATE ADVICE","Gemini 2.5 Flash\ncover letter\ninterview tips",GREEN),
]
for i,(n,head,body,col) in enumerate(nodes):
    x=.70+i*3.13
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,x,2.50,2.66,2.84,"131C32","2B385B",True,1)
    icon_badge(slide,n,x+.24,2.80,col,.48,9.5)
    add_text(slide,head,x+.24,3.53,2.15,.38,12.3,WHITE,True)
    add_text(slide,body,x+.24,4.13,2.1,.85,10.3,"A9B3C9")
    if i<3:
        add_text(slide,"→",x+2.70,3.64,.40,.28,21,"586888",True,align=PP_ALIGN.CENTER)

shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,.70,5.75,11.96,.61,"111B31","2B385B",True,1)
add_text(slide,"STATE",.96,5.96,.65,.18,8.5,"8EA0C8",True)
add_text(slide,"Inputs and outputs accumulate across nodes; errors stop safely, and missing results skip advice.",1.68,5.91,9.90,.24,11,"D3D9E7")
footer(slide, 7, dark=True)


# 8 — Architecture
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, WHITE, WHITE)
capabl_logo(slide)
title_block(slide, "System architecture", "Decoupled product, data and intelligence layers",
            "A versioned REST API connects the React client to persistent state and external services.")

# left client
shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,.67,2.25,2.18,3.78,"EEF2FF","C7D2FE",True,1)
add_text(slide,"CLIENT",.96,2.56,1.55,.18,8.5,INDIGO,True)
add_text(slide,"React 18",.96,3.00,1.55,.25,16,INK,True)
add_text(slide,"Vite + Tailwind",.96,3.36,1.55,.22,10.5,MUTED)
for i,t in enumerate(["Dashboard","Resume","Job search","Saved jobs","Interview"]):
    pill(slide,t,.96,3.87+i*.36,1.38,fill=WHITE,color=SLATE,size=8.2)

# api
shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,3.35,2.25,3.32,3.78,NAVY,NAVY,True,0)
add_text(slide,"APPLICATION API",3.70,2.56,2.0,.18,8.5,"8EA0C8",True)
add_text(slide,"FastAPI",3.70,3.00,1.60,.25,17,WHITE,True)
pill(slide,"/api/v1",5.38,2.98,.90,fill="202B48",color="A5B4FC",size=8.5)
api_items=[("AUTH","JWT + bcrypt"),("RESUME","parse + persist"),("MATCH","score + advice"),("JOBS","search + board"),("INTERVIEW","question + evaluate")]
for i,(a,b) in enumerate(api_items):
    add_text(slide,a,3.72,3.66+i*.41,.78,.18,8.2,"8290AA",True)
    add_text(slide,b,4.56,3.64+i*.41,1.68,.20,9.5,"D5DCEB")

# data
shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,7.18,2.25,2.24,3.78,PALE,LINE,True,1)
add_text(slide,"STATE",7.50,2.56,1.55,.18,8.5,SLATE,True)
for i,(head,body,col) in enumerate([("PostgreSQL","users • resumes • jobs",BLUE),("Redis","search cache • 1 hour",RED)]):
    y=3.08+i*1.22
    icon_badge(slide,"DB" if i==0 else "R",7.50,y,col,.44,9)
    add_text(slide,head,8.08,y-.01,1.06,.21,12,INK,True)
    add_text(slide,body,7.50,y+.48,1.55,.42,9,MUTED)

# external
shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,9.92,2.25,2.72,3.78,"F8FAFC",LINE,True,1)
add_text(slide,"EXTERNAL SERVICES",10.24,2.56,1.95,.18,8.5,SLATE,True)
for i,(head,body,col) in enumerate([("Adzuna","job listings",CYAN),("Gemini 2.5 Flash","advice + interviews",VIOLET)]):
    y=3.12+i*1.32
    icon_badge(slide,"A" if i==0 else "G",10.24,y,col,.48,11)
    add_text(slide,head,10.87,y,1.45,.24,11.5,INK,True)
    add_text(slide,body,10.87,y+.34,1.36,.35,9.2,MUTED)

for x1,x2 in [(2.85,3.35),(6.67,7.18),(9.42,9.92)]:
    line(slide,x1,4.05,x2,4.05,"94A3B8",2)
    add_text(slide,"→",x1+.12,3.84,.28,.22,14,"64748B",True,align=PP_ALIGN.CENTER)
footer(slide, 8)


# 9 — Engineering proof
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, PALE, PALE)
capabl_logo(slide)
title_block(slide, "Engineering snapshot", "What is built—and what gates production",
            "The current codebase is feature-complete for local demonstration, with explicit hardening work documented.")

cards=[
    ("53", "backend tests passing", "Agent, auth, interview, jobs, matching, resume parsing and users.", GREEN),
    ("19", "versioned API routes", "FastAPI endpoints across agent, auth, interview, jobs, match, resume and users.", BLUE),
    ("5", "product workspaces", "Dashboard, resume, search, saved jobs and mock interview.", VIOLET),
]
for i,(metric,head,body,col) in enumerate(cards):
    x=.68+i*4.13
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,x,2.20,3.78,1.78,WHITE,LINE,True,1)
    add_text(slide,metric,x+.25,2.46,.85,.45,27,col,True)
    add_text(slide,head,x+1.13,2.49,2.30,.28,12.5,INK,True)
    add_text(slide,body,x+.25,3.15,3.22,.47,9.7,MUTED)

shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,.68,4.30,7.24,2.08,WHITE,LINE,True,1)
add_text(slide,"DEPLOYMENT PATH",.98,4.62,1.85,.18,8.5,SLATE,True)
add_text(slide,"Frontend",.98,5.03,.92,.21,11.3,INK,True)
add_text(slide,"Vercel  •  npm build",1.92,5.03,2.0,.21,10.3,MUTED)
add_text(slide,"Backend",.98,5.42,.92,.21,11.3,INK,True)
add_text(slide,"Docker  •  Railway config",1.92,5.42,2.5,.21,10.3,MUTED)
add_text(slide,"Local",.98,5.81,.92,.21,11.3,INK,True)
add_text(slide,"Docker Compose  •  PostgreSQL  •  Redis",1.92,5.81,3.8,.21,10.3,MUTED)
shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,8.18,4.30,4.46,2.08,"FFF7ED","FED7AA",True,1)
add_text(slide,"PRODUCTION GATES",8.48,4.62,1.85,.18,8.5,"9A5A00",True)
for i,t in enumerate(["Enforce ownership checks on data routes", "Require a strong secret + restrict CORS", "Align backend host config and migrate Gemini SDK"]):
    add_bullet(slide,t,8.48,5.02+i*.40,3.70,dot=AMBER,size=9.9)
footer(slide, 9)


# 10 — Roadmap & close
slide = prs.slides.add_slide(blank)
rect(slide, 0, 0, W, H, NAVY, NAVY)
capabl_logo(slide, dark=True)
title_block(slide, "Roadmap + close", "A strong foundation for a smarter career loop",
            "The product already connects context across the journey; the next work improves trust, intelligence and reach.", dark=True)

phases=[
    ("NOW", "Demo-ready", ["Resume intelligence", "Ranked discovery", "Application board", "AI mock interviews"], BLUE),
    ("NEXT", "Production-safe", ["Route ownership", "Secret + CORS policy", "SDK migration", "Deployment alignment"], AMBER),
    ("LATER", "Learning system", ["Richer embeddings", "Feedback-driven ranking", "Notifications", "Outcome analytics"], GREEN),
]
for i,(tag,head,items,col) in enumerate(phases):
    x=.68+i*4.13
    shape(slide,MSO_SHAPE.ROUNDED_RECTANGLE,x,2.34,3.76,3.37,"131C32","2B385B",True,1)
    pill(slide,tag,x+.25,2.68,.72,fill="202B48",color=col,size=8.4)
    add_text(slide,head,x+.25,3.25,2.9,.30,17,WHITE,True)
    for j,t in enumerate(items):
        add_bullet(slide,t,x+.25,3.87+j*.42,3.05,color="B6C0D4",dot=col,size=10.2)

add_rich_text(slide,[
    ("CareerAI turns a resume from a document into ",16,False,"B9C3D8"),
    ("shared context for every career decision.",16,True,WHITE),
],.69,6.11,8.3,.38)
pill(slide,"THANK YOU",11.06,6.05,1.55,fill=INDIGO,color=WHITE,size=10)
footer(slide, 10, dark=True)


prs.save(OUT)
print(OUT)
